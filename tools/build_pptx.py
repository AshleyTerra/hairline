"""
Builds the Hairline Salon Manager presentation as a real, editable PowerPoint deck.

Design follows the Hairline wordmark: warm taupe on near-white paper, ink text,
generous whitespace, hairline rules. 16:9.
"""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

SHOTS = Path(
    r"C:\temp\claude\c--Data-OneDrive---Terra-Group-Applications-Hairline"
    r"\6305bd13-c82c-4c05-b471-9115eecd7529\scratchpad\guide-shots"
)
OUT = Path(r"C:\tmp\hairline-proto\docs\Hairline Salon Manager - Presentation.pptx")

# ----------------------------------------------------------------- brand
TAUPE = RGBColor(0x8A, 0x7F, 0x6F)
TAUPE_DEEP = RGBColor(0x6E, 0x64, 0x55)
TAUPE_SOFT = RGBColor(0xB3, 0xA8, 0x98)
INK = RGBColor(0x1A, 0x18, 0x16)
BODY = RGBColor(0x3A, 0x36, 0x2F)
MUTED = RGBColor(0x7A, 0x72, 0x64)
PAPER = RGBColor(0xFA, 0xF9, 0xF7)
CARD = RGBColor(0xFF, 0xFF, 0xFF)
RULE = RGBColor(0xE2, 0xDE, 0xD7)
CHIP = RGBColor(0xF1, 0xEE, 0xE8)
GOOD = RGBColor(0x4C, 0x7A, 0x5A)
WARN = RGBColor(0xA8, 0x76, 0x2A)
CRIT = RGBColor(0xA0, 0x43, 0x3A)
TEAL = RGBColor(0x0A, 0x86, 0xA8)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

FONT = "Segoe UI"
FONT_LIGHT = "Segoe UI Light"

W = Inches(13.333)
H = Inches(7.5)
M = Inches(0.9)  # page margin

prs = Presentation()
prs.slide_width = W
prs.slide_height = H
BLANK = prs.slide_layouts[6]


# ----------------------------------------------------------------- helpers
def slide(dark=False):
    s = prs.slides.add_slide(BLANK)
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = INK if dark else PAPER
    bg.line.fill.background()
    bg.shadow.inherit = False
    return s


def textbox(s, x, y, w, h, anchor=MSO_ANCHOR.TOP):
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    return tf


def para(tf, text, size, color, bold=False, font=FONT, space_after=0,
         align=PP_ALIGN.LEFT, spacing=None, first=False, caps=False, char_space=None):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    if spacing:
        p.line_spacing = spacing
    p.space_after = Pt(space_after)
    run = p.add_run()
    run.text = text.upper() if caps else text
    f = run.font
    f.name = font
    f.size = Pt(size)
    f.bold = bold
    f.color.rgb = color
    if char_space is not None:
        # Letter-spacing is not exposed by python-pptx; set it on the XML run.
        rPr = run._r.get_or_add_rPr()
        rPr.set("spc", str(int(char_space * 100)))
    return p


def rule(s, x, y, w, color=RULE):
    ln = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, Pt(0.75))
    ln.fill.solid()
    ln.fill.fore_color.rgb = color
    ln.line.fill.background()
    ln.shadow.inherit = False
    return ln


def wordmark(s, x, y, size=20, dark=False):
    """HAIR|line, taupe + ink, as on every other Hairline surface."""
    tf = textbox(s, x, y, Inches(3), Inches(0.6))
    p = tf.paragraphs[0]
    for text, colour, weight in (
        ("HAIR", TAUPE, FONT_LIGHT),
        ("|", WHITE if dark else INK, FONT_LIGHT),
        ("line", WHITE if dark else INK, FONT_LIGHT),
    ):
        run = p.add_run()
        run.text = text
        run.font.name = weight
        run.font.size = Pt(size)
        run.font.color.rgb = colour
    return tf


def eyebrow(s, text, y=None, dark=False):
    tf = textbox(s, M, y or Inches(1.15), Inches(9), Inches(0.3))
    para(tf, text, 10.5, TAUPE_SOFT if dark else TAUPE, bold=True, caps=True,
         char_space=1.8, first=True)


def heading(s, text, y=None, size=34, dark=False, width=Inches(10.5)):
    tf = textbox(s, M, y or Inches(1.5), width, Inches(1.3))
    para(tf, text, size, WHITE if dark else INK, font=FONT_LIGHT, spacing=1.05, first=True)


def card(s, x, y, w, h, fill=CARD, border=RULE):
    box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    box.fill.solid()
    box.fill.fore_color.rgb = fill
    box.line.color.rgb = border
    box.line.width = Pt(0.75)
    box.shadow.inherit = False
    box.adjustments[0] = 0.03
    return box


def stat(s, x, y, w, value, label, colour=INK, h=Inches(1.15)):
    card(s, x, y, w, h)
    tf = textbox(s, x + Inches(0.22), y + Inches(0.16), w - Inches(0.44), h - Inches(0.3))
    para(tf, value, 26, colour, font=FONT_LIGHT, first=True, space_after=2)
    para(tf, label, 9, MUTED, caps=True, char_space=0.8)


def slide_number(s, n, dark=False):
    tf = textbox(s, W - Inches(1.3), H - Inches(0.62), Inches(0.7), Inches(0.3))
    para(tf, f"{n:02d}", 9.5, MUTED if not dark else RGBColor(0x6B, 0x63, 0x57),
         align=PP_ALIGN.RIGHT, first=True)


def picture(s, name, x, y, w):
    """Places a screenshot with a hairline frame."""
    path = SHOTS / f"{name}.png"
    pic = s.shapes.add_picture(str(path), x, y, width=w)
    frame = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, pic.width, pic.height)
    frame.fill.background()
    frame.line.color.rgb = RULE
    frame.line.width = Pt(0.75)
    frame.shadow.inherit = False
    return pic


def bullets(s, x, y, w, items, size=13.5, dark=False, gap=9):
    """Dash-led list, matching the plan document and the web deck."""
    tf = textbox(s, x, y, w, Inches(3.6))
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap)
        p.line_spacing = 1.28
        dash = p.add_run()
        dash.text = "—   "
        dash.font.name = FONT
        dash.font.size = Pt(size)
        dash.font.color.rgb = TAUPE_SOFT if dark else TAUPE
        # A leading **bold** segment lets one phrase carry the line.
        if "**" in item:
            head, rest = item.split("**", 2)[1], item.split("**", 2)[2]
            r1 = p.add_run()
            r1.text = head
            r1.font.name = FONT
            r1.font.size = Pt(size)
            r1.font.bold = True
            r1.font.color.rgb = WHITE if dark else INK
            r2 = p.add_run()
            r2.text = rest
            r2.font.name = FONT
            r2.font.size = Pt(size)
            r2.font.color.rgb = RGBColor(0xCF, 0xC9, 0xBD) if dark else BODY
        else:
            r = p.add_run()
            r.text = item
            r.font.name = FONT
            r.font.size = Pt(size)
            r.font.color.rgb = RGBColor(0xCF, 0xC9, 0xBD) if dark else BODY
    return tf


# ================================================================= 1 title
s = slide()
tf = textbox(s, M, Inches(2.3), Inches(9), Inches(1.2))
p = tf.paragraphs[0]
for text, colour in (("HAIR", TAUPE), ("|", INK), ("line", INK)):
    r = p.add_run()
    r.text = text
    r.font.name = FONT_LIGHT
    r.font.size = Pt(66)
    r.font.color.rgb = colour

rule(s, M, Inches(3.5), Inches(7.2))

tf = textbox(s, M, Inches(3.85), Inches(9), Inches(1.6))
para(tf, "Salon Manager", 40, INK, font=FONT_LIGHT, first=True, space_after=10)
para(tf, "A working prototype, built on eleven years of Hairline's own trading data.",
     15, MUTED)

for i, (label, x) in enumerate([("Discovery", 0), ("Plan", 1.45), ("Live prototype", 2.6)]):
    w = Inches(1.3 if i < 2 else 1.75)
    chip = card(s, M + Inches(x), Inches(5.35), w, Inches(0.4), fill=CHIP, border=CHIP)
    tfc = textbox(s, M + Inches(x), Inches(5.44), w, Inches(0.3))
    para(tfc, label, 10, TAUPE_DEEP, bold=True, align=PP_ALIGN.CENTER, first=True)

tf = textbox(s, M, H - Inches(1.0), Inches(9), Inches(0.4))
para(tf, "Prepared for the owner of Hairline  Â·  Shop 30, Stoneridge Centre, Greenstone Park",
     10, MUTED, first=True)
slide_number(s, 1)

# ================================================================= 2 brief
s = slide()
wordmark(s, M, Inches(0.5), 15)
eyebrow(s, "The brief")
heading(s, "Replace MySalon with something the team will actually use.")

bullets(s, M, Inches(2.75), Inches(6.0), [
    "**Cover the whole salon** — but keep it simple. No overkill.",
    "**Focus on the core:** invoicing, client history, pricing, stock, staff.",
    "**Win buy-in** from reception, the stylists and the owner.",
    "**Build it properly** — this could become a product for other salons.",
])

card(s, Inches(7.6), Inches(2.6), Inches(4.85), Inches(2.5))
tf = textbox(s, Inches(7.95), Inches(2.9), Inches(4.15), Inches(2.0))
para(tf, "The test we set ourselves", 9.5, TAUPE, bold=True, caps=True, char_space=1.2,
     first=True, space_after=8)
para(tf, "The new till must be faster than the old one.", 17, INK, font=FONT_LIGHT,
     space_after=8, spacing=1.1)
para(tf, "Every study of salon software adoption says the same thing: if checkout slows "
         "down, the team stops using it — whatever else the system can do.", 11.5, MUTED,
     spacing=1.3)
slide_number(s, 2)

# ================================================================= 3 evidence
s = slide(dark=True)
wordmark(s, M, Inches(0.5), 15, dark=True)
eyebrow(s, "Before designing anything", dark=True)
heading(s, "We restored the salon's database\nand read eleven years of it.", dark=True)

tf = textbox(s, M, Inches(3.1), Inches(9.2), Inches(0.8))
para(tf, "Not a workshop, not assumptions — the actual record of how Hairline trades, "
         "taken from the July 2026 MySalon backup.", 14, RGBColor(0xCF, 0xC9, 0xBD),
     first=True, spacing=1.35)

data = [("90,183", "invoices since 2015"), ("8,645", "active clients"),
        ("R7.5m", "revenue in 2025"), ("92", "database tables read")]
x = M
for value, label in data:
    box = card(s, x, Inches(4.5), Inches(2.65), Inches(1.2),
               fill=RGBColor(0x24, 0x21, 0x1D), border=RGBColor(0x35, 0x30, 0x2A))
    tf = textbox(s, x + Inches(0.25), Inches(4.66), Inches(2.2), Inches(0.95))
    para(tf, value, 26, WHITE, font=FONT_LIGHT, first=True, space_after=2)
    para(tf, label, 9, RGBColor(0x8F, 0x87, 0x7A), caps=True, char_space=0.8)
    x += Inches(2.85)
slide_number(s, 3, dark=True)

# ================================================================= 4 used / unused
s = slide()
wordmark(s, M, Inches(0.5), 15)
eyebrow(s, "What the data said")
heading(s, "The salon told us what to build — and what to leave out.")

used = [("Invoicing", "90,183 invoices"), ("Cash-up", "3,841 daily counts"),
        ("Stock takes", "49,663 items counted"), ("Clocking in", "11,069 records"),
        ("Tips & advances", "4,478 records")]
unused = [("Payroll", "0 records"), ("Loyalty points", "0 records"),
          ("Packages", "0 records"), ("Bank reconciliation", "0 records"),
          ("Appointment diary", "48 in 11 years")]

for idx, (title, rows, colour) in enumerate(
    [("Used every single day", used, GOOD), ("Never used in 11 years", unused, CRIT)]
):
    x = M + Inches(idx * 5.9)
    tf = textbox(s, x, Inches(2.7), Inches(5.3), Inches(0.3))
    para(tf, title, 13, colour, bold=True, first=True)
    rule(s, x, Inches(3.05), Inches(5.3), colour)
    y = Inches(3.2)
    for name, value in rows:
        tfr = textbox(s, x, y, Inches(3.2), Inches(0.32))
        para(tfr, name, 12, INK, first=True)
        tfv = textbox(s, x + Inches(3.2), y, Inches(2.1), Inches(0.32))
        para(tfv, value, 11.5, MUTED, align=PP_ALIGN.RIGHT, first=True)
        y += Inches(0.46)
        rule(s, x, y - Inches(0.09), Inches(5.3), RGBColor(0xED, 0xEA, 0xE4))

tf = textbox(s, M, Inches(5.9), Inches(11.5), Inches(0.4))
para(tf, "Every unused feature we leave out is a screen reception never has to learn.",
     12, MUTED, first=True)
slide_number(s, 4)

# ================================================================= 5 findings
s = slide()
wordmark(s, M, Inches(0.5), 15)
eyebrow(s, "Four things the data revealed")
heading(s, "Findings the salon can act on today.")

findings = [
    ("01 — STOCK", "1,350 lines show\nnegative stock", CRIT,
     "Sales rung up against stock never received. Nearly 40% of the stock file cannot "
     "be trusted until it is counted."),
    ("02 — CLIENTS", "758 clients have\nquietly lapsed", WARN,
     "They visited within the year but not in 90 days. The cheapest revenue in the "
     "salon is a message to this list."),
    ("03 — PAYMENTS", "89% of takings\nare on card", TEAL,
     "Cash is barely 5%. The cash-up ritual guards a small and shrinking share of "
     "the money."),
    ("04 — MARKETING", "Only 8% have a\nbirthday on file", WARN,
     "718 of 8,645. Just 76 have an e-mail address, but 8,603 have a phone number — "
     "so SMS and WhatsApp, never e-mail."),
]
x = M
for num, title, colour, text in findings:
    card(s, x, Inches(2.7), Inches(2.65), Inches(3.4))
    tf = textbox(s, x + Inches(0.22), Inches(2.92), Inches(2.2), Inches(3.0))
    para(tf, num, 8.5, TAUPE, bold=True, char_space=1.0, first=True, space_after=7)
    para(tf, title, 15, colour, bold=True, spacing=1.12, space_after=9)
    para(tf, text, 10.5, MUTED, spacing=1.32)
    x += Inches(2.85)
slide_number(s, 5)

# ================================================================= 6 decision
s = slide(dark=True)
wordmark(s, M, Inches(0.5), 15, dark=True)
eyebrow(s, "The decision", dark=True)

tf = textbox(s, M, Inches(2.0), Inches(10.5), Inches(2.0))
p = tf.paragraphs[0]
p.line_spacing = 1.12
for text, colour in (("Build the ", WHITE), ("daily essentials", TAUPE_SOFT),
                     (" exceptionally well.", WHITE)):
    r = p.add_run()
    r.text = text
    r.font.name = FONT_LIGHT
    r.font.size = Pt(36)
    r.font.color.rgb = colour
p2 = tf.add_paragraph()
p2.line_spacing = 1.12
r = p2.add_run()
r.text = "Leave out everything the salon has already voted against."
r.font.name = FONT_LIGHT
r.font.size = Pt(36)
r.font.color.rgb = WHITE

rule(s, M, Inches(4.3), Inches(11.5), RGBColor(0x35, 0x30, 0x2A))

modules = ["Invoicing & cash-up", "Clients & history", "Pricing & menus", "Stock control",
           "Staff portfolios", "Client messaging", "A light diary"]
x = M
y = Inches(4.75)
for mod in modules:
    w = Inches(0.16 + len(mod) * 0.088)
    if x + w > W - M:
        x = M
        y += Inches(0.62)
    chip = card(s, x, y, w, Inches(0.44), fill=RGBColor(0x26, 0x22, 0x1D),
                border=RGBColor(0x35, 0x30, 0x2A))
    tfc = textbox(s, x, y + Inches(0.11), w, Inches(0.3))
    para(tfc, mod, 11, TAUPE_SOFT, align=PP_ALIGN.CENTER, first=True)
    x += w + Inches(0.16)
slide_number(s, 6, dark=True)

# ================================================================= 7 then / now
s = slide()
wordmark(s, M, Inches(0.5), 15)
eyebrow(s, "What changes")
heading(s, "Same rituals, twenty years of technology later.", size=30)

then_items = [
    "One PC in the salon. If it dies, trading stops.",
    "SQL Server 2005 under the hood.",
    "Price menus retyped in Word every year.",
    "Figures only visible at reception.",
    "Dated SMS; no WhatsApp.",
]
now_items = [
    "Cloud-hosted, with a till that keeps working offline.",
    "A modern web app on any device.",
    "Client menu generated from live prices.",
    "Takings on the owner's phone, anywhere.",
    "SMS now, WhatsApp next.",
]

# Left: how it works today
card(s, M, Inches(2.6), Inches(5.4), Inches(3.6))
tf = textbox(s, M + Inches(0.32), Inches(2.85), Inches(4.7), Inches(0.3))
para(tf, "MySalon today", 10, MUTED, bold=True, caps=True, char_space=1.4, first=True)
bullets(s, M + Inches(0.32), Inches(3.28), Inches(4.7), then_items, size=12, gap=10)

# Arrow
tfa = textbox(s, Inches(6.5), Inches(4.15), Inches(0.6), Inches(0.5))
para(tfa, "â†’", 26, TAUPE, align=PP_ALIGN.CENTER, first=True)

# Right: how it works next
box = card(s, Inches(7.1), Inches(2.6), Inches(5.35), Inches(3.6), fill=INK, border=INK)
tf = textbox(s, Inches(7.42), Inches(2.85), Inches(4.7), Inches(0.3))
para(tf, "Salon Manager", 10, TAUPE_SOFT, bold=True, caps=True, char_space=1.4, first=True)
bullets(s, Inches(7.42), Inches(3.28), Inches(4.7), now_items, size=12, dark=True, gap=10)
slide_number(s, 7)

# ================================================================= 8 prototype
s = slide(dark=True)
wordmark(s, M, Inches(0.5), 15, dark=True)
eyebrow(s, "Not a mock-up", dark=True)

tf = textbox(s, M, Inches(1.55), Inches(6.0), Inches(1.4))
p = tf.paragraphs[0]
p.line_spacing = 1.1
for text, colour in (("We built it. It's ", WHITE), ("running", TAUPE_SOFT),
                     (", on the salon's own data.", WHITE)):
    r = p.add_run()
    r.text = text
    r.font.name = FONT_LIGHT
    r.font.size = Pt(30)
    r.font.color.rgb = colour

tf = textbox(s, M, Inches(3.2), Inches(5.6), Inches(1.4))
para(tf, "Eight working screens loaded with real services, real prices, real stock and "
         "eleven years of real visit patterns. Client names are anonymised; nothing else "
         "is invented.", 13, RGBColor(0xCF, 0xC9, 0xBD), first=True, spacing=1.35)

for i, (value, label) in enumerate([("750", "client histories"), ("228", "services priced"),
                                    ("2,494", "stock lines")]):
    x = M + Inches(i * 1.95)
    tf = textbox(s, x, Inches(4.9), Inches(1.8), Inches(0.9))
    para(tf, value, 24, WHITE, font=FONT_LIGHT, first=True, space_after=2)
    para(tf, label, 9, RGBColor(0x8F, 0x87, 0x7A), caps=True, char_space=0.8)

picture(s, "02-dashboard", Inches(7.0), Inches(1.5), Inches(5.4))
slide_number(s, 8, dark=True)

# ================================================================= 9 the till
s = slide()
wordmark(s, M, Inches(0.5), 15)
eyebrow(s, "The screen that decides everything")
heading(s, "The till, with a stopwatch on it.", size=30)

picture(s, "05-till-payment", M, Inches(2.5), Inches(7.3))

bullets(s, Inches(8.7), Inches(2.6), Inches(3.8), [
    "**One screen:** client, services, products, discount, payment.",
    "**Split payments** across cash, card, EFT and vouchers.",
    "**Tips per stylist,** captured at the moment of payment.",
    "**A live timer** makes the 30-second target visible.",
], size=12, gap=11)

card(s, Inches(8.7), Inches(5.3), Inches(3.8), Inches(1.35), fill=CHIP, border=CHIP)
tf = textbox(s, Inches(8.95), Inches(5.5), Inches(3.3), Inches(1.0))
para(tf, "18 automated tests on the money logic", 11.5, TAUPE_DEEP, bold=True,
     first=True, space_after=5)
para(tf, "VAT, discounts, split payments, change and tips are each proven "
         "independently of the screen.", 10, TAUPE_DEEP, spacing=1.28)
slide_number(s, 9)

# ================================================================= 10 client history
s = slide()
wordmark(s, M, Inches(0.5), 15)
eyebrow(s, "The salon's most valuable asset")
heading(s, "Every client, every visit, instantly.", size=30)

picture(s, "07-client-file", Inches(5.0), Inches(2.5), Inches(7.4))

bullets(s, M, Inches(2.6), Inches(3.7), [
    "**Search** by name or phone; open a full timeline in one click.",
    "**Every service and price** back to 2015, with the stylist who did it.",
    "**Colour formulas** and medical flags travel with the client.",
    "**Filters** for lapsed clients, top spenders and birthdays.",
], size=12, gap=11)

card(s, M, Inches(5.35), Inches(3.7), Inches(1.3), fill=CHIP, border=CHIP)
tf = textbox(s, M + Inches(0.25), Inches(5.55), Inches(3.2), Inches(1.0))
para(tf, "1,709 clients have visited 10+ times", 11.5, TAUPE_DEEP, bold=True,
     first=True, space_after=5)
para(tf, "That loyal core is the business — and its history is locked in a 2005 "
         "database on one PC.", 10, TAUPE_DEEP, spacing=1.28)
slide_number(s, 10)

# ================================================================= 11 three views
s = slide()
wordmark(s, M, Inches(0.5), 15)
eyebrow(s, "Buy-in by design")
heading(s, "Three people, three different systems.", size=30)

views = [
    ("OWNER", "The business, on your phone",
     "Takings so far today, month against last, top stylists, stock value and what needs "
     "attention — without phoning the salon.", "02-dashboard"),
    ("RECEPTION", "A faster front desk",
     "Till, clients, diary, stock and the familiar denomination cash-up. No profit "
     "reports, no staff pay: less to learn.", "13-cashup"),
    ("STYLIST", "Their own numbers",
     "Today's clients, month against target and tips earned — on their own phone. Done "
     "for the team, not to them.", "16-stylist-dashboard"),
]
x = M
for label, title, text, shot in views:
    card(s, x, Inches(2.5), Inches(3.72), Inches(4.15))
    picture(s, shot, x + Inches(0.18), Inches(2.68), Inches(3.36))
    tf = textbox(s, x + Inches(0.25), Inches(4.75), Inches(3.2), Inches(1.7))
    para(tf, label, 8.5, TAUPE, bold=True, char_space=1.2, first=True, space_after=6)
    para(tf, title, 14, INK, bold=True, space_after=7)
    para(tf, text, 10.5, MUTED, spacing=1.3)
    x += Inches(3.95)
slide_number(s, 11)

# ================================================================= 12 rollout
s = slide()
wordmark(s, M, Inches(0.5), 15)
eyebrow(s, "Going live without risk")
heading(s, "Nothing switches over until the numbers match.", size=30)

phases = [
    ("Build the heart", "Till, cash-up and clients, loaded with migrated data", "6–8 weeks"),
    ("Shadow run", "The new till runs alongside MySalon; totals must match to the cent",
     "2 weeks"),
    ("Go live", "Invoicing and cash-up cut over; MySalon becomes read-only", "1 weekend"),
    ("Stock & team", "Stock control, orders, stock takes, portfolios, time clock",
     "+4–6 weeks"),
    ("Messaging & diary", "SMS and WhatsApp templates, the light diary, stylist views",
     "+3–4 weeks"),
]
y = Inches(2.65)
for i, (title, text, dur) in enumerate(phases):
    dot = s.shapes.add_shape(MSO_SHAPE.OVAL, M, y + Inches(0.09), Inches(0.16), Inches(0.16))
    dot.fill.solid()
    dot.fill.fore_color.rgb = PAPER
    dot.line.color.rgb = TAUPE
    dot.line.width = Pt(1.5)
    dot.shadow.inherit = False
    if i < len(phases) - 1:
        conn = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, M + Inches(0.072), y + Inches(0.25),
                                  Pt(0.75), Inches(0.55))
        conn.fill.solid()
        conn.fill.fore_color.rgb = RULE
        conn.line.fill.background()
        conn.shadow.inherit = False
    tf = textbox(s, M + Inches(0.45), y, Inches(8.0), Inches(0.6))
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title + "  "
    r.font.name = FONT
    r.font.size = Pt(13.5)
    r.font.bold = True
    r.font.color.rgb = INK
    r2 = p.add_run()
    r2.text = "— " + text
    r2.font.name = FONT
    r2.font.size = Pt(12.5)
    r2.font.color.rgb = BODY
    tfd = textbox(s, Inches(10.6), y, Inches(1.9), Inches(0.4))
    para(tfd, dur, 11, MUTED, align=PP_ALIGN.RIGHT, first=True)
    y += Inches(0.8)

tf = textbox(s, M, Inches(6.75), Inches(11.5), Inches(0.4))
para(tf, "Stylist dashboards arrive only after reception already trusts the till.",
     11.5, MUTED, first=True)
slide_number(s, 12)

# ================================================================= 13 success
s = slide(dark=True)
wordmark(s, M, Inches(0.5), 15, dark=True)
eyebrow(s, "How we'll know it worked", dark=True)
heading(s, "Five things, all of them measurable.", dark=True, size=30)

bullets(s, M, Inches(2.9), Inches(10.5), [
    "**A routine sale in 30 seconds or less,** and cash-up done in five minutes.",
    "**Shadow-run totals match MySalon to the cent,** two weeks running.",
    "**Reception say they prefer the new till** within a month of go-live.",
    "**The owner checks takings on their phone,** instead of phoning the salon.",
    "**Zero data loss** — every historic visit still findable.",
], size=14, dark=True, gap=14)
slide_number(s, 13, dark=True)

# ================================================================= 14 decisions
s = slide()
wordmark(s, M, Inches(0.5), 15)
eyebrow(s, "Over to you")
heading(s, "Try the prototype, then five decisions.", size=30)

card(s, M, Inches(2.55), Inches(5.4), Inches(3.1), fill=CHIP, border=CHIP)
tf = textbox(s, M + Inches(0.3), Inches(2.78), Inches(4.8), Inches(2.6))
para(tf, "What to do in the prototype", 9.5, TAUPE_DEEP, bold=True, caps=True,
     char_space=1.2, first=True, space_after=10)
for line in ["Sign in as owner — password hairline2026.",
             "Ring up a sale on the Till and watch the timer.",
             "Open a client and scroll eleven years of history.",
             "Sign out, back in as karin, to see a stylist's view.",
             "Count the drawer on Cash-up."]:
    p = tf.add_paragraph()
    p.space_after = Pt(7)
    p.line_spacing = 1.25
    d = p.add_run()
    d.text = "—   "
    d.font.name = FONT
    d.font.size = Pt(12)
    d.font.color.rgb = TAUPE
    r = p.add_run()
    r.text = line
    r.font.name = FONT
    r.font.size = Pt(12)
    r.font.color.rgb = BODY

x2 = Inches(7.0)
tf = textbox(s, x2, Inches(2.55), Inches(5.4), Inches(0.3))
para(tf, "Decisions we need", 9.5, TAUPE, bold=True, caps=True, char_space=1.2, first=True)
rule(s, x2, Inches(2.85), Inches(5.4), TAUPE)
y = Inches(3.0)
for name, question in [("Name", "Keep \u201cHairline Salon Manager\u201d?"),
                       ("Price tiers", "Confirm senior / junior levels"),
                       ("Card machine", "Standalone, or integrated later?"),
                       ("WhatsApp", "Register the salon number?"),
                       ("Diary", "From day one, or ease into it?")]:
    tfn = textbox(s, x2, y, Inches(1.9), Inches(0.34))
    para(tfn, name, 12, INK, bold=True, first=True)
    tfq = textbox(s, x2 + Inches(1.9), y, Inches(3.5), Inches(0.34))
    para(tfq, question, 11.5, MUTED, first=True)
    y += Inches(0.52)
    rule(s, x2, y - Inches(0.1), Inches(5.4), RGBColor(0xED, 0xEA, 0xE4))

tf = textbox(s, M, Inches(6.5), Inches(11.5), Inches(0.4))
para(tf, "This plan is a conversation starter, not a contract. Everything in it can change.",
     11.5, MUTED, first=True)
slide_number(s, 14)

OUT.parent.mkdir(parents=True, exist_ok=True)
prs.save(str(OUT))
print(f"Saved {OUT}  ({OUT.stat().st_size / 1024:.0f} KB, {len(prs.slides.__iter__.__self__._sldIdLst)} slides)")
